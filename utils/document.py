import os 


def extract_text(file_path):
	ext = os.path.splitext(file_path)[1].lower()
	text = ""
	if ext == ".pdf":
		import fitz  # PyMuPDF
		with fitz.open(file_path) as doc:
			# text = text + doc[1].get_text()
			for page in doc: text += page.get_text()
	elif ext == ".docx":
		from docx import Document
		doc = Document(file_path)
		text = "\n".join([p.text for p in doc.paragraphs])
	elif ext == ".pptx":
		from pptx import Presentation
		prs = Presentation(file_path)
		for slide in prs.slides:
			for shape in slide.shapes:
				if hasattr(shape, "text"): text += shape.text + " "
	else:
		# For other file types, try reading as plain text
		try:
			with open(file_path, 'r', encoding='utf-8') as f:
				text = f.read()
		except Exception as e:
			print(f"Could not read file {file_path}: {e}")
	return text
